import os
from pptx import Presentation
from pptx.util import Inches, Pt


def generate_ppt_sync(outline_id: str, images: list, template: str, options: dict) -> str:
    """
    DEV stub: generate a simple PPTX using python-pptx.
    In production, outline_id should be resolved to actual content.
    """
    prs = Presentation()
    # Simple slide generation for demo
    for i in range(3):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        body = slide.shapes.placeholders[1]
        title.text = f"示範投影片 {i+1}"
        tf = body.text_frame
        tf.text = "- 範例重點 1"
        p = tf.add_paragraph()
        p.text = "- 範例重點 2"
        p.level = 1

    out_dir = os.environ.get("BACKEND_STORAGE_PATH", "/workspaces/ppt/backend/storage")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, f"presentation_{abs(hash(outline_id)) % 100000}.pptx")
    prs.save(out_path)
    return out_path
